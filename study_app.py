import streamlit as st
import openai
import json
import io
import base64
from typing import List, Dict, Any
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import re

# Import configuration
try:
    from config import OPENAI_API_KEY
    # Check if API key is still the placeholder
    if OPENAI_API_KEY == "sk-your-actual-openai-api-key-here":
        OPENAI_API_KEY = None
except ImportError:
    # Fallback - no API key found
    OPENAI_API_KEY = None

# Language configurations
LANGUAGES = {
    "English": {
        "code": "en",
        "name": "English",
        "translations": {
            "title": "📚 AI Study Assistant",
            "subtitle": "Upload your study materials and get AI-powered summaries and quizzes!",
            "configuration": "🔑 Configuration",
            "language_setting": "🌐 Language",
            "quiz_progress": "📊 Quiz Progress",
            "upload_material": "📤 Upload Material",
            "summary": "📋 Summary",
            "quiz": "🧠 Quiz", 
            "results": "📈 Results",
            "upload_header": "Upload Your Study Materials",
            "choose_file": "Choose a file (PDF, Word, PowerPoint, or Text)",
            "file_help": "Upload lecture slides, notes, or any study material",
            "paste_text": "Or paste your text directly:",
            "paste_placeholder": "Copy and paste your notes, lecture content, or study material here...",
            "process_material": "🚀 Process Material",
            "processing": "Processing your study material...",
            "generating_summary": "Generating summary...",
            "creating_quiz": "Creating quiz questions...",
            "success_message": "✅ Material processed successfully! Check the Summary and Quiz tabs.",
            "upload_warning": "Please upload a file or paste text to continue.",
            "summary_header": "📋 Study Material Summary",
            "summary_info": "Upload and process your study material first to see the summary.",
            "quiz_header": "🧠 Interactive Quiz",
            "quiz_info": "Upload and process your study material first to start the quiz.",
            "quiz_completed": "🎉 Quiz completed! Check the Results tab to see your performance.",
            "question_of": "Question {current} of {total}",
            "choose_answer": "Choose your answer:",
            "previous": "⬅️ Previous",
            "submit_answer": "✅ Submit Answer",
            "next": "➡️ Next",
            "finish_quiz": "🏁 Finish Quiz",
            "correct": "✅ Correct! Well done!",
            "incorrect": "❌ Incorrect. Let's review this concept.",
            "key_concept": "💡 Key Concept:",
            "explanation": "Explanation:",
            "correct_answer": "Correct Answer:",
            "results_header": "📈 Quiz Results",
            "results_info": "Complete the quiz to see your results.",
            "total_questions": "Total Questions",
            "correct_answers": "Correct Answers",
            "score": "Score",
            "excellent": "🌟 Excellent! You have a great understanding of the material!",
            "good": "👍 Good job! You understand most of the concepts.",
            "fair": "📚 Not bad, but consider reviewing the material again.",
            "poor": "📖 You might want to study the material more thoroughly.",
            "detailed_review": "📝 Detailed Review",
            "concepts_to_review": "🎯 Concepts to Review:",
            "question_review": "🔍 Question by Question Review",
            "retake_quiz": "🔄 Take Quiz Again"
        }
    },
    "Japanese": {
        "code": "ja",
        "name": "日本語",
        "translations": {
            "title": "📚 AI学習アシスタント",
            "subtitle": "学習資料をアップロードして、AI搭載の要約とクイズを取得しましょう！",
            "configuration": "🔑 設定",
            "language_setting": "🌐 言語",
            "quiz_progress": "📊 クイズ進捗",
            "upload_material": "📤 資料アップロード",
            "summary": "📋 要約",
            "quiz": "🧠 クイズ",
            "results": "📈 結果",
            "upload_header": "学習資料をアップロード",
            "choose_file": "ファイルを選択（PDF、Word、PowerPoint、またはテキスト）",
            "file_help": "講義スライド、ノート、または学習資料をアップロード",
            "paste_text": "またはテキストを直接貼り付け：",
            "paste_placeholder": "ノート、講義内容、または学習資料をここにコピー＆ペーストしてください...",
            "process_material": "🚀 資料を処理",
            "processing": "学習資料を処理中...",
            "generating_summary": "要約を生成中...",
            "creating_quiz": "クイズ問題を作成中...",
            "success_message": "✅ 資料の処理が完了しました！要約とクイズタブをご確認ください。",
            "upload_warning": "続行するにはファイルをアップロードまたはテキストを貼り付けてください。",
            "summary_header": "📋 学習資料要約",
            "summary_info": "要約を見るには、まず学習資料をアップロードして処理してください。",
            "quiz_header": "🧠 インタラクティブクイズ",
            "quiz_info": "クイズを開始するには、まず学習資料をアップロードして処理してください。",
            "quiz_completed": "🎉 クイズ完了！結果タブでパフォーマンスを確認してください。",
            "question_of": "問題 {current} / {total}",
            "choose_answer": "答えを選択してください：",
            "previous": "⬅️ 前へ",
            "submit_answer": "✅ 答えを送信",
            "next": "➡️ 次へ",
            "finish_quiz": "🏁 クイズ終了",
            "correct": "✅ 正解！よくできました！",
            "incorrect": "❌ 不正解。この概念を復習しましょう。",
            "key_concept": "💡 重要概念：",
            "explanation": "説明：",
            "correct_answer": "正解：",
            "results_header": "📈 クイズ結果",
            "results_info": "結果を見るにはクイズを完了してください。",
            "total_questions": "総問題数",
            "correct_answers": "正解数",
            "score": "スコア",
            "excellent": "🌟 素晴らしい！教材をよく理解しています！",
            "good": "👍 よくできました！ほとんどの概念を理解しています。",
            "fair": "📚 悪くありませんが、もう一度復習することをお勧めします。",
            "poor": "📖 教材をもっと徹底的に勉強した方がよいかもしれません。",
            "detailed_review": "📝 詳細レビュー",
            "concepts_to_review": "🎯 復習すべき概念：",
            "question_review": "🔍 問題別レビュー",
            "retake_quiz": "🔄 クイズを再受験"
        }
    },
    "Korean": {
        "code": "ko",
        "name": "한국어",
        "translations": {
            "title": "📚 AI 학습 도우미",
            "subtitle": "학습 자료를 업로드하고 AI 기반 요약 및 퀴즈를 받아보세요!",
            "configuration": "🔑 설정",
            "language_setting": "🌐 언어",
            "quiz_progress": "📊 퀴즈 진행률",
            "upload_material": "📤 자료 업로드",
            "summary": "📋 요약",
            "quiz": "🧠 퀴즈",
            "results": "📈 결과",
            "upload_header": "학습 자료 업로드",
            "choose_file": "파일 선택 (PDF, Word, PowerPoint 또는 텍스트)",
            "file_help": "강의 슬라이드, 노트 또는 학습 자료 업로드",
            "paste_text": "또는 텍스트를 직접 붙여넣기:",
            "paste_placeholder": "노트, 강의 내용 또는 학습 자료를 여기에 복사하여 붙여넣으세요...",
            "process_material": "🚀 자료 처리",
            "processing": "학습 자료 처리 중...",
            "generating_summary": "요약 생성 중...",
            "creating_quiz": "퀴즈 문제 생성 중...",
            "success_message": "✅ 자료 처리가 완료되었습니다! 요약 및 퀴즈 탭을 확인하세요.",
            "upload_warning": "계속하려면 파일을 업로드하거나 텍스트를 붙여넣으세요.",
            "summary_header": "📋 학습 자료 요약",
            "summary_info": "요약을 보려면 먼저 학습 자료를 업로드하고 처리하세요.",
            "quiz_header": "🧠 인터랙티브 퀴즈",
            "quiz_info": "퀴즈를 시작하려면 먼저 학습 자료를 업로드하고 처리하세요.",
            "quiz_completed": "🎉 퀴즈 완료! 결과 탭에서 성과를 확인하세요.",
            "question_of": "문제 {current} / {total}",
            "choose_answer": "답을 선택하세요:",
            "previous": "⬅️ 이전",
            "submit_answer": "✅ 답 제출",
            "next": "➡️ 다음",
            "finish_quiz": "🏁 퀴즈 완료",
            "correct": "✅ 정답! 잘했습니다!",
            "incorrect": "❌ 오답입니다. 이 개념을 복습해봅시다.",
            "key_concept": "💡 핵심 개념:",
            "explanation": "설명:",
            "correct_answer": "정답:",
            "results_header": "📈 퀴즈 결과",
            "results_info": "결과를 보려면 퀴즈를 완료하세요.",
            "total_questions": "총 문제수",
            "correct_answers": "정답수",
            "score": "점수",
            "excellent": "🌟 훌륭합니다! 자료를 잘 이해하고 있습니다!",
            "good": "👍 잘했습니다! 대부분의 개념을 이해하고 있습니다.",
            "fair": "📚 나쁘지 않지만 자료를 다시 검토해보는 것을 고려해보세요.",
            "poor": "📖 자료를 더 철저히 공부하는 것이 좋겠습니다.",
            "detailed_review": "📝 상세 리뷰",
            "concepts_to_review": "🎯 복습할 개념:",
            "question_review": "🔍 문제별 리뷰",
            "retake_quiz": "🔄 퀴즈 다시 풀기"
        }
    },
    "Spanish": {
        "code": "es",
        "name": "Español",
        "translations": {
            "title": "📚 Asistente de Estudio IA",
            "subtitle": "¡Sube tus materiales de estudio y obtén resúmenes y cuestionarios con IA!",
            "configuration": "🔑 Configuración",
            "language_setting": "🌐 Idioma",
            "quiz_progress": "📊 Progreso del Cuestionario",
            "upload_material": "📤 Subir Material",
            "summary": "📋 Resumen",
            "quiz": "🧠 Cuestionario",
            "results": "📈 Resultados",
            "upload_header": "Sube tus Materiales de Estudio",
            "choose_file": "Elige un archivo (PDF, Word, PowerPoint o Texto)",
            "file_help": "Sube diapositivas de clase, notas o cualquier material de estudio",
            "paste_text": "O pega tu texto directamente:",
            "paste_placeholder": "Copia y pega tus notas, contenido de clase o material de estudio aquí...",
            "process_material": "🚀 Procesar Material",
            "processing": "Procesando tu material de estudio...",
            "generating_summary": "Generando resumen...",
            "creating_quiz": "Creando preguntas del cuestionario...",
            "success_message": "✅ ¡Material procesado exitosamente! Revisa las pestañas de Resumen y Cuestionario.",
            "upload_warning": "Por favor sube un archivo o pega texto para continuar.",
            "summary_header": "📋 Resumen del Material de Estudio",
            "summary_info": "Sube y procesa tu material de estudio primero para ver el resumen.",
            "quiz_header": "🧠 Cuestionario Interactivo",
            "quiz_info": "Sube y procesa tu material de estudio primero para comenzar el cuestionario.",
            "quiz_completed": "🎉 ¡Cuestionario completado! Revisa la pestaña de Resultados para ver tu rendimiento.",
            "question_of": "Pregunta {current} de {total}",
            "choose_answer": "Elige tu respuesta:",
            "previous": "⬅️ Anterior",
            "submit_answer": "✅ Enviar Respuesta",
            "next": "➡️ Siguiente",
            "finish_quiz": "🏁 Terminar Cuestionario",
            "correct": "✅ ¡Correcto! ¡Bien hecho!",
            "incorrect": "❌ Incorrecto. Repasemos este concepto.",
            "key_concept": "💡 Concepto Clave:",
            "explanation": "Explicación:",
            "correct_answer": "Respuesta Correcta:",
            "results_header": "📈 Resultados del Cuestionario",
            "results_info": "Completa el cuestionario para ver tus resultados.",
            "total_questions": "Total de Preguntas",
            "correct_answers": "Respuestas Correctas",
            "score": "Puntuación",
            "excellent": "🌟 ¡Excelente! ¡Tienes un gran entendimiento del material!",
            "good": "👍 ¡Buen trabajo! Entiendes la mayoría de los conceptos.",
            "fair": "📚 No está mal, pero considera repasar el material de nuevo.",
            "poor": "📖 Tal vez quieras estudiar el material más a fondo.",
            "detailed_review": "📝 Revisión Detallada",
            "concepts_to_review": "🎯 Conceptos para Repasar:",
            "question_review": "🔍 Revisión Pregunta por Pregunta",
            "retake_quiz": "🔄 Tomar Cuestionario de Nuevo"
        }
    },
    "French": {
        "code": "fr",
        "name": "Français",
        "translations": {
            "title": "📚 Assistant d'Étude IA",
            "subtitle": "Téléchargez vos supports d'étude et obtenez des résumés et quiz alimentés par l'IA !",
            "configuration": "🔑 Configuration",
            "language_setting": "🌐 Langue",
            "quiz_progress": "📊 Progression du Quiz",
            "upload_material": "📤 Télécharger du Matériel",
            "summary": "📋 Résumé",
            "quiz": "🧠 Quiz",
            "results": "📈 Résultats",
            "upload_header": "Téléchargez vos Supports d'Étude",
            "choose_file": "Choisir un fichier (PDF, Word, PowerPoint ou Texte)",
            "file_help": "Téléchargez des diapositives de cours, notes ou tout matériel d'étude",
            "paste_text": "Ou collez votre texte directement :",
            "paste_placeholder": "Copiez et collez vos notes, contenu de cours ou matériel d'étude ici...",
            "process_material": "🚀 Traiter le Matériel",
            "processing": "Traitement de votre matériel d'étude...",
            "generating_summary": "Génération du résumé...",
            "creating_quiz": "Création des questions du quiz...",
            "success_message": "✅ Matériel traité avec succès ! Vérifiez les onglets Résumé et Quiz.",
            "upload_warning": "Veuillez télécharger un fichier ou coller du texte pour continuer.",
            "summary_header": "📋 Résumé du Matériel d'Étude",
            "summary_info": "Téléchargez et traitez d'abord votre matériel d'étude pour voir le résumé.",
            "quiz_header": "🧠 Quiz Interactif",
            "quiz_info": "Téléchargez et traitez d'abord votre matériel d'étude pour commencer le quiz.",
            "quiz_completed": "🎉 Quiz terminé ! Vérifiez l'onglet Résultats pour voir votre performance.",
            "question_of": "Question {current} sur {total}",
            "choose_answer": "Choisissez votre réponse :",
            "previous": "⬅️ Précédent",
            "submit_answer": "✅ Soumettre la Réponse",
            "next": "➡️ Suivant",
            "finish_quiz": "🏁 Terminer le Quiz",
            "correct": "✅ Correct ! Bien joué !",
            "incorrect": "❌ Incorrect. Révisons ce concept.",
            "key_concept": "💡 Concept Clé :",
            "explanation": "Explication :",
            "correct_answer": "Réponse Correcte :",
            "results_header": "📈 Résultats du Quiz",
            "results_info": "Terminez le quiz pour voir vos résultats.",
            "total_questions": "Total des Questions",
            "correct_answers": "Réponses Correctes",
            "score": "Score",
            "excellent": "🌟 Excellent ! Vous avez une excellente compréhension du matériel !",
            "good": "👍 Bon travail ! Vous comprenez la plupart des concepts.",
            "fair": "📚 Pas mal, mais envisagez de réviser le matériel à nouveau.",
            "poor": "📖 Vous pourriez vouloir étudier le matériel plus en profondeur.",
            "detailed_review": "📝 Révision Détaillée",
            "concepts_to_review": "🎯 Concepts à Réviser :",
            "question_review": "🔍 Révision Question par Question",
            "retake_quiz": "🔄 Refaire le Quiz"
        }
    },
    "Chinese": {
        "code": "zh",
        "name": "中文",
        "translations": {
            "title": "📚 AI学习助手",
            "subtitle": "上传您的学习材料，获得AI驱动的摘要和测验！",
            "configuration": "🔑 配置",
            "language_setting": "🌐 语言",
            "quiz_progress": "📊 测验进度",
            "upload_material": "📤 上传材料",
            "summary": "📋 摘要",
            "quiz": "🧠 测验",
            "results": "📈 结果",
            "upload_header": "上传您的学习材料",
            "choose_file": "选择文件（PDF、Word、PowerPoint或文本）",
            "file_help": "上传讲座幻灯片、笔记或任何学习材料",
            "paste_text": "或直接粘贴您的文本：",
            "paste_placeholder": "在此复制粘贴您的笔记、讲座内容或学习材料...",
            "process_material": "🚀 处理材料",
            "processing": "正在处理您的学习材料...",
            "generating_summary": "正在生成摘要...",
            "creating_quiz": "正在创建测验问题...",
            "success_message": "✅ 材料处理成功！请查看摘要和测验选项卡。",
            "upload_warning": "请上传文件或粘贴文本以继续。",
            "summary_header": "📋 学习材料摘要",
            "summary_info": "请先上传和处理您的学习材料以查看摘要。",
            "quiz_header": "🧠 互动测验",
            "quiz_info": "请先上传和处理您的学习材料以开始测验。",
            "quiz_completed": "🎉 测验完成！请查看结果选项卡以查看您的表现。",
            "question_of": "第 {current} 题，共 {total} 题",
            "choose_answer": "选择您的答案：",
            "previous": "⬅️ 上一题",
            "submit_answer": "✅ 提交答案",
            "next": "➡️ 下一题",
            "finish_quiz": "🏁 完成测验",
            "correct": "✅ 正确！做得很好！",
            "incorrect": "❌ 错误。让我们复习这个概念。",
            "key_concept": "💡 关键概念：",
            "explanation": "解释：",
            "correct_answer": "正确答案：",
            "results_header": "📈 测验结果",
            "results_info": "完成测验以查看您的结果。",
            "total_questions": "总题数",
            "correct_answers": "正确答案数",
            "score": "得分",
            "excellent": "🌟 优秀！您对材料有很好的理解！",
            "good": "👍 做得好！您理解了大部分概念。",
            "fair": "📚 不错，但建议再次复习材料。",
            "poor": "📖 您可能需要更彻底地学习材料。",
            "detailed_review": "📝 详细回顾",
            "concepts_to_review": "🎯 需要复习的概念：",
            "question_review": "🔍 逐题回顾",
            "retake_quiz": "🔄 重新参加测验"
        }
    }
}

# Configure page
st.set_page_config(
    page_title="AI Study Assistant",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for modern UI
st.markdown("""
<style>
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        text-align: center;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 2rem;
    }
    
    .quiz-question {
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
        border-left: 4px solid #667eea;
        background: #f8f9ff;
    }
    
    .summary-content {
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
        background: #f8f9ff;
        border-left: 4px solid #667eea;
    }
    
    .correct-answer {
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        color: #155724;
        padding: 1rem;
        border-radius: 10px;
        margin: 0.5rem 0;
    }
    
    .incorrect-answer {
        background-color: #f8d7da;
        border: 1px solid #f5c6cb;
        color: #721c24;
        padding: 1rem;
        border-radius: 10px;
        margin: 0.5rem 0;
    }
    
    .key-concept {
        background-color: #fff3cd;
        border: 1px solid #ffeaa7;
        color: #856404;
        padding: 1rem;
        border-radius: 10px;
        margin: 0.5rem 0;
        font-weight: bold;
    }
    
    .progress-bar {
        background-color: #e9ecef;
        border-radius: 10px;
        height: 20px;
        margin: 1rem 0;
    }
    
    /* Language selector styling */
    div[data-testid="column"]:last-child .stSelectbox > div > div {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        border-radius: 20px;
        border: none;
        color: white;
        font-weight: bold;
        min-height: 35px;
    }
    
    div[data-testid="column"]:last-child .stSelectbox > div > div > div {
        color: white;
    }
    
    /* Hide the arrow for language selector */
    div[data-testid="column"]:last-child .stSelectbox svg {
        fill: white;
    }
    
    /* File uploader styling */
    .stFileUploader > div > div > div > div {
        border: 2px dashed #667eea !important;
        border-radius: 15px !important;
        background: linear-gradient(135deg, #f8f9ff 0%, #e6f2ff 100%) !important;
        padding: 2rem !important;
        text-align: center !important;
    }
    
    .stFileUploader > div > div > div > div > div {
        font-size: 1.1rem !important;
        color: #667eea !important;
        font-weight: 500 !important;
    }
    
    /* Enhanced drag and drop area */
    .stFileUploader > div > div > div > div:hover {
        background: linear-gradient(135deg, #e6f2ff 0%, #d4e9ff 100%) !important;
        border-color: #4a6cf7 !important;
        transition: all 0.3s ease !important;
    }
</style>
""", unsafe_allow_html=True)

class StudyMaterialProcessor:
    """Process various types of study materials"""
    
    @staticmethod
    def extract_text_from_pdf(file) -> str:
        """Extract text from PDF file"""
        try:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PDF: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(file) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading Word document: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PowerPoint: {str(e)}")
            return ""

class AIQuizGenerator:
    """Generate quizzes and summaries using OpenAI API"""
    
    def __init__(self, api_key: str, language_code: str = "en"):
        # Check if using OpenRouter or similar service
        if api_key.startswith("sk-or-"):
            # Configure for OpenRouter
            self.client = openai.OpenAI(
                api_key=api_key,
                base_url="https://openrouter.ai/api/v1"
            )
        else:
            # Standard OpenAI configuration
            self.client = openai.OpenAI(api_key=api_key)
        self.language_code = language_code
    
    def _get_language_instruction(self) -> str:
        """Get language instruction for OpenAI based on selected language"""
        language_instructions = {
            "en": "Respond in English.",
            "ja": "日本語で回答してください。",
            "ko": "한국어로 답변해주세요.",
            "es": "Responde en español.",
            "fr": "Répondez en français.",
            "zh": "请用中文回答。"
        }
        return language_instructions.get(self.language_code, "Respond in English.")
    
    def generate_summary(self, text: str) -> str:
        """Generate a summary of the study material"""
        try:
            language_instruction = self._get_language_instruction()
            response = self.client.chat.completions.create(
                model="openai/gpt-oss-120b:free",
                messages=[
                    {"role": "system", "content": f"You are an expert educational assistant. Create a comprehensive summary of the provided study material, highlighting the most important concepts, key points, and main ideas. Structure the summary with clear headings and bullet points. {language_instruction}"},
                    {"role": "user", "content": f"Please summarize this study material:\n\n{text[:4000]}"}
                ],
                max_tokens=1000,
                temperature=0.3
            )
            return response.choices[0].message.content
        except Exception as e:
            error_msg = str(e)
            if "401" in error_msg or "api_key" in error_msg.lower():
                st.error("🔑 API Key Error: Please check your API key is valid and has sufficient credits.")
            elif "429" in error_msg:
                st.error("⏰ Rate Limit: Please wait a moment and try again.")
            else:
                st.error(f"Error generating summary: {error_msg}")
            return "Unable to generate summary due to API error. Please check your API key and try again."
    
    def generate_quiz(self, text: str) -> List[Dict[str, Any]]:
        """Generate 10 quiz questions based on the material"""
        try:
            language_instruction = self._get_language_instruction()
            response = self.client.chat.completions.create(
                model="openai/gpt-oss-120b:free",
                messages=[
                    {"role": "system", "content": f"""You are an expert quiz creator. Generate exactly 10 multiple-choice questions based on the provided study material. {language_instruction}
                    
                    IMPORTANT: Return ONLY a valid JSON array. Do not include any explanatory text, markdown formatting, or code blocks.
                    
                    Each question must have this exact format:
                    {{
                        "question": "Question text here",
                        "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
                        "correct_answer": "A",
                        "explanation": "Detailed explanation of why this is correct",
                        "key_concept": "The main concept this question tests"
                    }}
                    
                    Return as a JSON array starting with [ and ending with ]. Make sure questions cover the most important concepts from the material."""},
                    {"role": "user", "content": f"Create 10 quiz questions from this material:\n\n{text[:3000]}"}
                ],
                max_tokens=2500,
                temperature=0.3
            )
            
            quiz_text = response.choices[0].message.content
            
            # Debug: Check if we got a response
            if not quiz_text or quiz_text.strip() == "":
                st.error("🤖 Empty Response: The AI service returned an empty response. Please try again.")
                return []
            
            # Clean up the response to ensure it's valid JSON
            quiz_text = quiz_text.strip()
            if quiz_text.startswith("```json"):
                quiz_text = quiz_text[7:]
            if quiz_text.startswith("```"):
                quiz_text = quiz_text[3:]
            if quiz_text.endswith("```"):
                quiz_text = quiz_text[:-3]
            
            # Remove any leading/trailing whitespace again
            quiz_text = quiz_text.strip()
            
            # Debug: Show first part of response if it's not valid JSON
            try:
                quiz_data = json.loads(quiz_text)
                if not isinstance(quiz_data, list) or len(quiz_data) == 0:
                    st.error("🔍 Invalid Format: Expected a list of questions but got a different format.")
                    return []
                return quiz_data
            except json.JSONDecodeError as json_error:
                st.error(f"📄 JSON Parse Error: The AI response wasn't valid JSON.")
                st.error(f"Response preview: {quiz_text[:200]}...")
                return []
                
        except Exception as e:
            error_msg = str(e)
            if "401" in error_msg or "api_key" in error_msg.lower():
                st.error("🔑 API Key Error: Please check your API key is valid and has sufficient credits.")
            elif "429" in error_msg:
                st.error("⏰ Rate Limit: Please wait a moment and try again.")
            elif "model" in error_msg.lower() and "not found" in error_msg.lower():
                st.error("🤖 Model Error: The specified AI model is not available. Please check your API service.")
            else:
                st.error(f"Error generating quiz: {error_msg}")
            
            # Fallback: Create a basic quiz about study techniques
            st.warning("🔄 Using fallback quiz about general study techniques.")
            return self._create_fallback_quiz()
    
    def _create_fallback_quiz(self) -> List[Dict[str, Any]]:
        """Create a fallback quiz when AI generation fails"""
        fallback_quiz = [
            {
                "question": "What is one of the most effective study techniques for long-term retention?",
                "options": ["A) Cramming the night before", "B) Spaced repetition", "C) Reading only once", "D) Memorizing without understanding"],
                "correct_answer": "B",
                "explanation": "Spaced repetition involves reviewing material at increasing intervals, which helps move information from short-term to long-term memory.",
                "key_concept": "Memory and Learning Techniques"
            },
            {
                "question": "Which method helps improve understanding of complex topics?",
                "options": ["A) Passive reading", "B) Active recall", "C) Highlighting everything", "D) Copying notes verbatim"],
                "correct_answer": "B", 
                "explanation": "Active recall involves testing yourself on the material, which strengthens neural pathways and improves comprehension.",
                "key_concept": "Active Learning Strategies"
            },
            {
                "question": "What is the recommended study session length for optimal focus?",
                "options": ["A) 10-15 minutes", "B) 25-50 minutes", "C) 2-3 hours", "D) 5-6 hours"],
                "correct_answer": "B",
                "explanation": "Research shows that 25-50 minute study sessions with breaks help maintain concentration and prevent mental fatigue.",
                "key_concept": "Time Management"
            }
        ]
        return fallback_quiz

def init_session_state():
    """Initialize session state variables"""
    if 'quiz_data' not in st.session_state:
        st.session_state.quiz_data = []
    if 'current_question' not in st.session_state:
        st.session_state.current_question = 0
    if 'user_answers' not in st.session_state:
        st.session_state.user_answers = {}
    if 'quiz_completed' not in st.session_state:
        st.session_state.quiz_completed = False
    if 'study_material' not in st.session_state:
        st.session_state.study_material = ""
    if 'summary' not in st.session_state:
        st.session_state.summary = ""
    if 'selected_language' not in st.session_state:
        st.session_state.selected_language = "English"

def setup_language_in_settings():
    """Set up language selection in a clean top-right interface"""
    # Create a language selector in the top right area using columns
    col1, col2, col3 = st.columns([6, 3, 3])
    
    with col3:
        # Create a language flag selector
        language_flags = {
            "English": "🇬🇧",
            "Japanese": "🇯🇵", 
            "Korean": "🇰🇷",
            "Spanish": "🇪🇸",
            "French": "🇫🇷",
            "Chinese": "🇨🇳"
        }
        
        current_flag = language_flags[st.session_state.selected_language]
        language_names = list(LANGUAGES.keys())
        current_index = language_names.index(st.session_state.selected_language)
        
        # Create a compact language selector
        selected_language = st.selectbox(
            "🌐",
            options=language_names,
            index=current_index,
            format_func=lambda x: f"{language_flags[x]} {LANGUAGES[x]['name']}",
            label_visibility="collapsed",
            key="top_language_selector"
        )
        
        # Update selected language
        if st.session_state.selected_language != selected_language:
            st.session_state.selected_language = selected_language
            st.rerun()

def get_text(key: str) -> str:
    """Get translated text based on selected language"""
    selected_lang = st.session_state.selected_language
    return LANGUAGES[selected_lang]["translations"].get(key, key)

def main():
    init_session_state()
    setup_language_in_settings()
    
    # Check for API key
    api_key = OPENAI_API_KEY
    if not api_key or api_key == "No API key found":
        st.error("🔑 API Key Required")
        st.info("Please add your API key to continue using the AI features.")
        st.markdown("""
        **How to get your API key:**
        
        **For OpenRouter (recommended for free models):**
        1. Visit [OpenRouter](https://openrouter.ai/)
        2. Sign up and get your API key (starts with `sk-or-`)
        3. Add it to the `config.py` file
        
        **For OpenAI:**
        1. Visit [OpenAI Platform](https://platform.openai.com/)
        2. Sign in or create an account
        3. Go to API Keys section
        4. Create a new API key (starts with `sk-`)
        5. Add it to the `config.py` file
        """)
        
        # Temporary API key input
        with st.sidebar:
            st.header("🔑 API Key Setup")
            temp_api_key = st.text_input(
                "Enter your API Key:",
                type="password",
                help="OpenRouter (sk-or-...) or OpenAI (sk-...) API key"
            )
            if temp_api_key:
                api_key = temp_api_key
                st.success("API key provided for this session!")
        
        if not api_key:
            return
    
    # Header
    st.markdown(f'<h1 class="main-header">{get_text("title")}</h1>', unsafe_allow_html=True)
    st.markdown(f'<p style="text-align: center; font-size: 1.2rem; color: #666;">{get_text("subtitle")}</p>', unsafe_allow_html=True)
    
    # Sidebar for progress tracking only
    with st.sidebar:
        # Progress tracker
        if st.session_state.quiz_data:
            st.header(get_text("quiz_progress"))
            progress = len(st.session_state.user_answers) / len(st.session_state.quiz_data)
            st.progress(progress)
            st.write(f"{get_text('question_of').format(current=len(st.session_state.user_answers), total=len(st.session_state.quiz_data))}")
        else:
            st.header(get_text("configuration"))
            st.info("Upload study material to start learning!")
    
    # Main content area
    tab1, tab2, tab3, tab4 = st.tabs([
        get_text("upload_material"), 
        get_text("summary"), 
        get_text("quiz"), 
        get_text("results")
    ])
    
    with tab1:
        st.header(get_text("upload_header"))
        
        # File upload area - simplified without custom styling box
        uploaded_file = st.file_uploader(
            get_text("choose_file"),
            type=['pdf', 'docx', 'pptx', 'txt'],
            help=get_text("file_help"),
            label_visibility="visible"
        )
        
        # Manual text input option
        st.subheader(get_text("paste_text"))
        manual_text = st.text_area(
            get_text("paste_text"), 
            height=200, 
            placeholder=get_text("paste_placeholder")
        )
        
        if st.button(get_text("process_material"), type="primary"):
            if uploaded_file or manual_text:
                with st.spinner(get_text("processing")):
                    # Process uploaded file
                    text_content = ""
                    if uploaded_file:
                        processor = StudyMaterialProcessor()
                        if uploaded_file.type == "application/pdf":
                            text_content = processor.extract_text_from_pdf(uploaded_file)
                        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                            text_content = processor.extract_text_from_docx(uploaded_file)
                        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                            text_content = processor.extract_text_from_pptx(uploaded_file)
                        elif uploaded_file.type == "text/plain":
                            text_content = str(uploaded_file.read(), "utf-8")
                    
                    # Use manual text if no file uploaded
                    if manual_text:
                        text_content = manual_text
                    
                    if text_content.strip():
                        st.session_state.study_material = text_content
                        
                        # Get current language code
                        current_lang_code = LANGUAGES[st.session_state.selected_language]["code"]
                        
                        # Generate summary and quiz
                        ai_generator = AIQuizGenerator(OPENAI_API_KEY, current_lang_code)
                        
                        with st.spinner(get_text("generating_summary")):
                            st.session_state.summary = ai_generator.generate_summary(text_content)
                        
                        with st.spinner(get_text("creating_quiz")):
                            st.session_state.quiz_data = ai_generator.generate_quiz(text_content)
                        
                        # Reset quiz state
                        st.session_state.current_question = 0
                        st.session_state.user_answers = {}
                        st.session_state.quiz_completed = False
                        
                        st.success(get_text("success_message"))
                        st.balloons()
                    else:
                        st.error("Could not extract text from the uploaded file. Please try a different file or paste text manually.")
            else:
                st.warning(get_text("upload_warning"))
    
    with tab2:
        st.header(get_text("summary_header"))
        if st.session_state.summary:
            st.markdown('<div class="summary-content">', unsafe_allow_html=True)
            st.markdown(st.session_state.summary)
            st.markdown('</div>', unsafe_allow_html=True)
        else:
            st.info(get_text("summary_info"))
    
    with tab3:
        st.header(get_text("quiz_header"))
        
        if not st.session_state.quiz_data:
            st.info(get_text("quiz_info"))
            return
        
        if st.session_state.quiz_completed:
            st.success(get_text("quiz_completed"))
            return
        
        # Display current question
        current_q = st.session_state.current_question
        if current_q < len(st.session_state.quiz_data):
            question_data = st.session_state.quiz_data[current_q]
            
            st.markdown('<div class="quiz-question">', unsafe_allow_html=True)
            st.subheader(get_text("question_of").format(current=current_q + 1, total=len(st.session_state.quiz_data)))
            st.write(question_data['question'])
            
            # Display options
            user_answer = st.radio(
                get_text("choose_answer"),
                options=question_data['options'],
                key=f"q_{current_q}",
                index=None
            )
            
            col1, col2, col3 = st.columns([1, 1, 1])
            
            with col1:
                if st.button(get_text("previous"), disabled=(current_q == 0)):
                    st.session_state.current_question = max(0, current_q - 1)
                    st.rerun()
            
            with col2:
                if st.button(get_text("submit_answer"), disabled=(user_answer is None)):
                    # Store answer
                    answer_letter = user_answer[0]  # Get A, B, C, or D
                    st.session_state.user_answers[current_q] = {
                        'answer': answer_letter,
                        'correct': answer_letter == question_data['correct_answer'],
                        'question_data': question_data
                    }
                    
                    # Show immediate feedback
                    if answer_letter == question_data['correct_answer']:
                        st.markdown(f'<div class="correct-answer">{get_text("correct")}</div>', unsafe_allow_html=True)
                    else:
                        st.markdown(f'<div class="incorrect-answer">{get_text("incorrect")}</div>', unsafe_allow_html=True)
                        st.markdown(f'<div class="key-concept">{get_text("key_concept")} {question_data["key_concept"]}</div>', unsafe_allow_html=True)
                        st.markdown(f"**{get_text('explanation')}** {question_data['explanation']}")
                    
                    st.write(f"**{get_text('correct_answer')}** {question_data['correct_answer']}) {question_data['options'][ord(question_data['correct_answer']) - ord('A')]}")
            
            with col3:
                if current_q in st.session_state.user_answers:
                    if current_q < len(st.session_state.quiz_data) - 1:
                        if st.button(get_text("next")):
                            st.session_state.current_question = current_q + 1
                            st.rerun()
                    else:
                        if st.button(get_text("finish_quiz")):
                            st.session_state.quiz_completed = True
                            st.rerun()
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    with tab4:
        st.header(get_text("results_header"))
        
        if not st.session_state.user_answers:
            st.info(get_text("results_info"))
            return
        
        # Calculate score
        total_questions = len(st.session_state.quiz_data)
        correct_answers = sum(1 for answer in st.session_state.user_answers.values() if answer['correct'])
        score_percentage = (correct_answers / total_questions) * 100
        
        # Display score
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric(get_text("total_questions"), total_questions)
        with col2:
            st.metric(get_text("correct_answers"), correct_answers)
        with col3:
            st.metric(get_text("score"), f"{score_percentage:.1f}%")
        
        # Score interpretation
        if score_percentage >= 90:
            st.success(get_text("excellent"))
        elif score_percentage >= 70:
            st.success(get_text("good"))
        elif score_percentage >= 50:
            st.warning(get_text("fair"))
        else:
            st.error(get_text("poor"))
        
        # Detailed results
        st.subheader(get_text("detailed_review"))
        
        # Show concepts that need review
        wrong_concepts = []
        for i, answer_data in st.session_state.user_answers.items():
            if not answer_data['correct']:
                wrong_concepts.append(answer_data['question_data']['key_concept'])
        
        if wrong_concepts:
            st.subheader(get_text("concepts_to_review"))
            for concept in set(wrong_concepts):  # Remove duplicates
                st.markdown(f'<div class="key-concept">💡 {concept}</div>', unsafe_allow_html=True)
        
        # Question by question review
        with st.expander(get_text("question_review")):
            for i in range(len(st.session_state.quiz_data)):
                if i in st.session_state.user_answers:
                    answer_data = st.session_state.user_answers[i]
                    question_data = answer_data['question_data']
                    
                    if answer_data['correct']:
                        st.markdown(f"**Q{i+1}:** ✅ Correct")
                    else:
                        st.markdown(f"**Q{i+1}:** ❌ Incorrect")
                        st.write(f"**Your answer:** {answer_data['answer']}")
                        st.write(f"**{get_text('correct_answer')}** {question_data['correct_answer']}")
                        st.write(f"**Key concept:** {question_data['key_concept']}")
                        st.write(f"**{get_text('explanation')}** {question_data['explanation']}")
                    st.write("---")
        
        # Reset quiz button
        if st.button(get_text("retake_quiz")):
            st.session_state.current_question = 0
            st.session_state.user_answers = {}
            st.session_state.quiz_completed = False
            st.rerun()

if __name__ == "__main__":
    main()
